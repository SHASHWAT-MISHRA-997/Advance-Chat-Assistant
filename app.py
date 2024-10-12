import streamlit as st 
import ollama
from PyPDF2 import PdfReader
from docx import Document
import pandas as pd
import tempfile
import os
import moviepy.editor as mp
from deep_translator import GoogleTranslator
import pptx
import requests
from bs4 import BeautifulSoup
import datetime
import speech_recognition as sr
import spacy
from gtts import gTTS
import threading
from transformers import pipeline
import base64



# ---------------------------
# Function Definitions
# ---------------------------

# Load SpaCy English model
try:
    nlp = spacy.load("en_core_web_sm")
except:
    st.error("SpaCy model 'en_core_web_sm' not found. Please run 'python -m spacy download en_core_web_sm' in your terminal.")
    st.stop()

# Initialize sentiment and emotion analysis pipeline
sentiment_pipeline = pipeline("sentiment-analysis")

emotion_pipeline = pipeline("text-classification", model="j-hartmann/emotion-english-distilroberta-base", return_all_scores=True)

# Initialize speech recognition and text-to-speech
recognizer = sr.Recognizer()
#engine = pyttsx3.init() 

# Optionally set a voice
voices = engine.getProperty('voices')
engine.setProperty('voice', voices[0].id)  # Change index for different voices

# Function to get response from Llama chatbot
def llama_chatbot(message, context=""):
    prompt = f"{context}\nUser: {message}\nBot:"
    response = ollama.generate(model='llama3.2:3b', prompt=prompt)
    return response["response"]

# Function to extract text from uploaded PDFs
def extract_pdf_text(uploaded_files):
    text = ""
    for uploaded_file in uploaded_files:
        try:
            reader_pdf = PdfReader(uploaded_file)
            for page in reader_pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        except Exception as e:
            st.error(f"Error reading {uploaded_file.name}: {e}")
    return text

# Function to extract text from uploaded Word files
def extract_word_text(uploaded_files):
    text = ""
    for uploaded_file in uploaded_files:
        try:
            doc = Document(uploaded_file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        except Exception as e:
            st.error(f"Error reading {uploaded_file.name}: {e}")
    return text

# Function to extract text from uploaded PPTX files
def extract_pptx_text(uploaded_files):
    text = ""
    for uploaded_file in uploaded_files:
        try:
            prs = pptx.Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            st.error(f"Error reading {uploaded_file.name}: {e}")
    return text

# Function to extract text from uploaded Excel files
def extract_excel_text(uploaded_files):
    text = ""
    for uploaded_file in uploaded_files:
        try:
            df = pd.read_excel(uploaded_file)
            text += df.to_string() + "\n"
        except Exception as e:
            st.error(f"Error reading {uploaded_file.name}: {e}")
    return text

# Function to extract text from a website URL
def extract_website_text(url):
    try:
        response = requests.get(url)
        soup = BeautifulSoup(response.text, 'html.parser')
        # Remove script and style elements
        for script_or_style in soup(["script", "style"]):
            script_or_style.decompose()
        text = ' '.join(soup.stripped_strings)
        return text
    except Exception as e:
        st.error(f"Error fetching {url}: {e}")
        return ""

# Function to extract audio from video and save it as a file
def extract_audio_from_video(video_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as tmp_video_file:
        tmp_video_file.write(video_file.read())
        tmp_video_file_path = tmp_video_file.name

    with tempfile.NamedTemporaryFile(delete=False, suffix=".aac") as tmp_audio_file:
        audio_file_path = tmp_audio_file.name

    video = mp.VideoFileClip(tmp_video_file_path)

    try:
        video.audio.write_audiofile(audio_file_path, codec='aac')
    finally:
        video.close()

    os.remove(tmp_video_file_path)

    return audio_file_path

# Function to translate text to a target language
def translate_text(text, dest_language):
    try:
        translated = GoogleTranslator(source='auto', target=dest_language).translate(text)
        return translated
    except Exception as e:
        st.error(f"Error translating text: {e}")
        return ""

# Function to listen to user input via microphone
def listen_to_user():
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        st.info("Listening... 🎤")
        try:
            audio = recognizer.listen(source, timeout=5)  # Wait for user to speak
            user_input = recognizer.recognize_google(audio)
            st.success("Audio recognized successfully! 🎉")
            return user_input
        except sr.WaitTimeoutError:
            st.error("Listening timed out while waiting for phrase to start.")
            return None
        except sr.UnknownValueError:
            st.error("Could not understand audio, please try again.")
            return None
        except sr.RequestError as e:
            st.error(f"Could not request results from Google Speech Recognition service; {e}")
            return None
        except Exception as e:
            st.error(f"An error occurred: {e}")
            return None

# Function to convert natural language to SQL using Llama chatbot
def text_to_sql(query):
    prompt = f"Convert the following natural language query into an SQL statement:\n\nQuery: {query}\nSQL:"
    response = ollama.generate(model='llama3.2:3b', prompt=prompt)
    sql_query = response.get("response", "")
    return sql_query


# Function to generate podcast audio from text
def generate_podcast_from_text(text):
    try:
        tts = gTTS(text)
        podcast_path = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3").name
        tts.save(podcast_path)
        return podcast_path
    except Exception as e:
        st.error(f"Error generating podcast: {e}")
        return None


# Function to extract text from uploaded audio files
def transcribe_audio(uploaded_audio):
    recognizer = sr.Recognizer()
    with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp_audio_file:
        tmp_audio_file.write(uploaded_audio.read())
        audio_path = tmp_audio_file.name

    try:
        with sr.AudioFile(audio_path) as source:
            audio = recognizer.record(source)
            text = recognizer.recognize_google(audio)
            return text
    except sr.UnknownValueError:
        st.error("Could not understand the audio.")
        return ""
    except sr.RequestError as e:
        st.error(f"Error with the speech recognition service: {e}")
        return ""
    except Exception as e:
        st.error(f"An error occurred: {e}")
        return ""


# Function to analyze sentiment of text
def analyze_sentiment(text):
    try:
        result = sentiment_pipeline(text)
        return result
    except Exception as e:
        st.error(f"Error analyzing sentiment: {e}")

# Function to analyze emotions in text
def analyze_emotions(text):
    try:
        result = emotion_pipeline(text)
        return result
    except Exception as e:
        st.error(f"Error analyzing emotions: {e}")


# Function to speak text using pyttsx3 in a separate thread
def speak_text(text):
    def speak():
        try:
            engine.say(text)
            if not engine._inLoop:
                engine.runAndWait()
        except RuntimeError as e:
            print(f"Error: {e}")
    
    # Run the speak function in a separate thread to avoid blocking
    thread = threading.Thread(target=speak)
    thread.start()


# Function to recognize speech using the microphone
def recognize_speech():
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        st.info("Listening for voice command... 🎤")
        try:
            audio = recognizer.listen(source, timeout=5)
            user_command = recognizer.recognize_google(audio)
            return user_command
        except sr.UnknownValueError:
            st.error("Could not understand the audio. Please try again.")
            return None
        except sr.RequestError as e:
            st.error(f"Error with the speech recognition service: {e}")
            return None

def set_background(image_file_path):
    global encoded_image  # Declare encoded_image as global to use it outside the function
    # Open the image file
    with open(image_file_path, "rb") as image_file:
        # Convert the image to base64 encoding
        encoded_image = base64.b64encode(image_file.read()).decode()

# ---------------------------
# Streamlit App Layout
# ---------------------------

# Initialize session state for tracking speaking status and input
if 'is_speaking' not in st.session_state:
    st.session_state['is_speaking'] = False
if 'user_input' not in st.session_state:
    st.session_state['user_input'] = ""


# Configure the Streamlit page
st.set_page_config(page_title="🤖 SM Business Chatbot", page_icon="💬", layout="wide")

# Initialize session state for chat history and user activity
if 'history' not in st.session_state:
    st.session_state['history'] = []
if 'activity' not in st.session_state:
    st.session_state['activity'] = []

with st.sidebar:
    # Set the background image
    set_background("ml.jpeg")
    
    st.markdown(
    f"""
    <style>
    /* Set the sidebar background color */
    .stSidebar {{
        background-color: rgba(0, 0, 0, 0.7) !important; /* Semi-transparent black for visibility */
        color: white !important; /* Sidebar text stays white */
    }}

    /* Ensure all elements in sidebar remain white */
    .stSidebar * {{
        color: white !important; /* Force all sidebar text to white */
    }}

    </style>
    """,
    unsafe_allow_html=True
)

    # Main content area with specific black font color
    st.markdown(
        f"""
        <style>
        /* Apply background, black text, and bold font to the main content area */
        .stApp {{
            background-image: url("data:image/jpeg;base64,{encoded_image}");
            background-size: cover;  /* Cover will adapt to screen size */
            background-position: center;
            background-attachment: fixed; /* Fixed background for a parallax effect */
            font-weight: bold; /* Set font to bold for all main content text */
        }}

        /* Set text color for the main content area to black */
        .stApp * {{
            color: white !important; /* Set all text in the main content to black */
        }}

        /* Prevent any sidebar text styling from affecting main content */
        .stSidebar * {{
            color: white !important; /* Ensure sidebar elements are white */
        }}
        </style>
        """,
        unsafe_allow_html=True
    )

    

# ---------------------------
    st.markdown("""
        <style>
        /* RGB Animation */
        @keyframes rgb {
            0% {background-position: 0% 50%;}
            50% {background-position: 100% 50%;}
            100% {background-position: 0% 50%;}
        }

        /* Background for sidebar and main area */
        .css-1y4d28q {
            background: linear-gradient(270deg, #ff0000, #00ff00, #0000ff);
            background-size: 600% 600%;
            animation: rgb 10s ease infinite;
            padding: 20px; /* Adding padding for better look */
        }

        /* Button styling with RGB effects */
        .stButton button {
            background: linear-gradient(270deg, #ff0000, #00ff00, #0000ff);
            background-size: 600% 600%;
            color: white;
            border: none;
            padding: 10px 20px;
            font-size: 16px;
            border-radius: 8px;
            animation: rgb 3s ease infinite;
            transition: all 0.3s ease;
        }

        /* Hover effect for buttons */
        .stButton button:hover {
            transform: scale(1.1);
            box-shadow: 0px 0px 15px rgba(255, 255, 255, 0.8);
            cursor: pointer;
        }

        /* Tab hover effect */
        .css-18e3th9 a {
            background: linear-gradient(270deg, #ff0000, #00ff00, #0000ff);
            background-size: 400% 400%;
            color: white;
            padding: 10px;
            border-radius: 10px;
            transition: 0.3s;
        }

        /* On hover, add glowing effect and scaling */
        .css-18e3th9 a:hover {
            animation: rgb 5s ease infinite;
            transform: scale(1.1);
            box-shadow: 0px 0px 20px rgba(255, 255, 255, 0.8);
        }

        /* RGB Animation for the top text (Created by SHASHWAT MISHRA) */
        .creator-link {
            text-align: center;
            font-size: 20px;
            margin-top: -50px;
            padding-bottom: 10px;
            color: white;
            background: linear-gradient(270deg, #ff0000, #00ff00, #0000ff);
            background-size: 600% 600%;
            animation: rgb 4s ease infinite;
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
        }
        </style>
    """, unsafe_allow_html=True)

    st.markdown("""
    <style>
    .creator-link {
        text-align: center;
        font-size: 20px;
        margin-top: -50px;
        padding-bottom: 10px;
        color: #007bff;
    }
    </style>
    <p class="creator-link">Created by <a href="https://www.linkedin.com/in/sm980/" target="_blank">SHASHWAT MISHRA</a></p>
    """, unsafe_allow_html=True)
    with st.sidebar:
        # Display the image in the sidebar
        st.image("ai.jpeg", use_column_width=True)  # Adjust the size to fit the sidebar

    st.title("🤖💼 SM Business & Research Assistant")
    st.markdown(""" 
    Welcome to the **SM Business & Research Assistant**! Empower your business operations and research activities by seamlessly interacting with documents, websites, videos, and more. 🚀

    ### 📌 Features:
    - **💬 Chat with AI-BOT:** Ask any questions and receive instant, AI-powered responses.
    - **📄 Chat with Documents:** Upload PDF, Word, PPTX, or Excel documents and get insights.
    - **🌐 Web Content Analyzer:** Input website URLs to fetch and interact with their content.
    - **📹 Video to Audio Extraction:** Upload video files, extract audio, and ask questions about the content.
    - **🎙️ Podcast Generator:** Turn your PDF documents into audio podcasts for easy consumption.
    - **🌍 Multilingual Support:** Translate and interact in multiple languages, including Hindi and other regional languages.
    - **📊 Natural Language to SQL Converter:** Convert natural language queries into structured SQL statements for database operations.
    - **🧑‍🏫 Research Assistant:** Upload research documents, analyze content, and ask targeted research questions.

    ### 🛠️ Instructions:
    1. **💬 Chat with AI-BOT:** Choose between voice or text input, type your message or speak, and receive AI-generated responses.
    2. **📄 Upload Documents:** Navigate to the **"📄 Chat with Documents"** section, upload your documents, and start chatting based on their content.
    3. **🌐 Input Website URLs:** In the **"🌐 Web Content Analyzer"** section, input the URL to fetch and interact with the website’s content.
    4. **📹 Upload Videos:** Go to the **"📹 Video to Audio Extraction"** section, upload your video, and ask questions based on the audio content.
    5. **🎙️ Podcast Generator:** Go to the **"🎙️ Podcast Generator"** section, upload a PDF to convert its content into a podcast.
    6. **🌍 Translate Text:** Use the translation tool to convert any text or document into your preferred language.
    7. **📊 Text to SQL:** Enter a natural language query and get the corresponding SQL statement.
    8. **🧑‍🏫 Research Assistant:** Upload research documents in the **"🧑‍🏫 Research Assistant"** section to extract insights and ask questions about the research material.

    Experience a streamlined, professional workflow for your business and research needs! 😊
    """)
# Streamlit App Layout
st.header("🗣️ Voice-Activated Controls")

# Option to view the Voice Command Guide
if st.checkbox("View Voice Command Guide 🛠️"):
    st.subheader("Voice Activated Controls Guide")
    st.markdown("""
    ## 🗣️ How to Use Voice Commands Across All Tabs

    ### **1. General Commands**
    These commands can be used across multiple tabs to perform common actions:
    
    - **"Upload file"**: Opens the file uploader so you can upload files (PDFs, DOCX, PPTX, Excel).
    - **"Analyze"**: Starts analyzing the content of an uploaded file or document.
    - **"Switch tab"**: Switches between different tabs in the app, e.g., moving from 'Chat with AI-BOT' to 'Chat with Documents'.
    - **"Translate to [language]"**: Translates text or documents to a target language (e.g., "Translate to Spanish").
    - **"Generate podcast"**: Converts an uploaded document (PDF) into an audio podcast.
    - **"Ask a question"**: Ask any general question (e.g., "What is artificial intelligence?") and get an AI response.
    
    ---

    ### **2. Tab-Specific Commands**

    #### 💬 **Tab: Chat with AI-BOT**
    - **Command**: "Ask [your question here]" (e.g., "What is the capital of France?")
    - **Task**: Chat with the AI-BOT and get answers or explanations to your questions.
    
    #### 📄 **Tab: Chat with Documents**
    - **Command**: "Upload file" followed by "Analyze" or ask a specific question related to the document content.
    - **Task**: Upload PDF, DOCX, PPTX, or Excel files and ask the AI to analyze or answer questions about the document’s content.

    #### 🌐 **Tab: Chat with Websites**
    - **Command**: "Fetch website" followed by a URL and "Analyze" or ask a question about the website’s content.
    - **Task**: Enter a website URL to fetch its content, then ask the AI-BOT to analyze the site or provide insights.

    #### 🎥🎤 **Tab: Chat with Media**
    - **Command**: "Upload video" or "Upload audio" followed by "Analyze" or "Transcribe".
    - **Task**: Upload a video to extract audio or upload an audio file to transcribe the speech into text. You can then ask the AI to analyze or provide insights into the content.

    #### 🎙️ **Tab: Podcast from PDF**
    - **Command**: "Generate podcast" after uploading a PDF.
    - **Task**: Converts a PDF document into an audio podcast for easy listening.

    #### 🌍 **Tab: AI Translator**
    - **Command**: "Translate to [language]" (e.g., "Translate to French").
    - **Task**: Translate input text or an uploaded document into the selected language.

    #### 📊 **Tab: Text to SQL**
    - **Command**: "Convert to SQL" followed by a natural language query (e.g., "Show all users where age is greater than 30").
    - **Task**: Converts your natural language query into an SQL statement.

    #### 🧑‍🏫 **Tab: Research Assistant**
    - **Command**: "Upload research document" followed by "Analyze" or ask specific research-related questions.
    - **Task**: Upload research papers or documents and ask questions related to the research topic. The AI will analyze and provide insights.

    #### 🧠 **Tab: Sentiment & Emotion Analysis**
    - **Command**: "Analyze sentiment" or "Analyze emotions".
    - **Task**: Analyze the sentiment or emotions of the uploaded text or document to understand the underlying tone.

    ---

    ### **Voice Command Tips**
    - **Speak Clearly**: Ensure you’re in a quiet environment to improve recognition accuracy.
    - **Command Structure**: Use simple, clear commands like "Upload file", "Analyze document", or "Switch tab".
    - **Wait for Confirmation**: After speaking, wait for the app to confirm the command it recognized before proceeding.

    """)

# Button to trigger voice recognition
if st.button("Speak 🎤", key="speak_button"):

    voice_input = recognize_speech()
    
    if voice_input:
        st.success(f"Recognized Command: {voice_input}")
        speak_text(f"You said: {voice_input}")  # Use the updated speak_text function

        # Example of simple voice command actions
        if "upload" in voice_input.lower():
            st.info("You can upload files by selecting the 'Upload' option.")
        elif "analyze" in voice_input.lower():
            st.info("Starting analysis of your document or data.")
        elif "switch tab" in voice_input.lower():
            st.info("Switching tab for you.")
        elif "translate" in voice_input.lower():
            st.info("Translating text into the requested language.")
        elif "generate podcast" in voice_input.lower():
            st.info("Generating a podcast from the document content.")
        elif "ask" in voice_input.lower():
            st.info("Ask a question for the AI to answer.")
        else:
            st.info("Command not recognized. Please try again or refer to the guide.")

# Tabs for different functionalities
tab1, tab2, tab3, tab4, tab5, tab6, tab7, tab8, tab9 = st.tabs(
    ["💬 Chat with AI-BOT", "📄 Chat with Documents", "🌐 Chat with Websites", 
     "🎥🎤 Chat with Media", "🎙️ Podcast from PDF", "🌍 AI Translator", 
     "📊 Text to SQL", "🧑‍🏫 Research Assistant", "🧠 Sentiment & Emotion Analysis"]
)

# ---------------------------
# Tab 1: Chat with AI-BOT
# ---------------------------
with tab1:
    st.title("🤖💬 Chat with AI Bot")

    # Display chat history only in this tab
    if st.session_state['history']:
        for chat in st.session_state['history']:
            st.markdown(f"**You:** {chat['user']} ({chat['timestamp']})")
            st.markdown(f"**Bot:** {chat['bot']}")

    st.markdown("---")

    # Initialize user input
    user_input = ""

    # Radio button to choose input method (Only visible in this tab)
    input_method = st.radio("Choose input method:", ("Text", "Voice"))

    if input_method == "Text":
        # Get user input without needing a button
        user_input = st.text_input("Type your message here and press Enter:", key="user_input")

        if user_input:
            if user_input.strip() != "":
                with st.spinner("Thinking... 🤔"):
                    bot_response = llama_chatbot(user_input)
                    timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    st.session_state['history'].append({
                        "user": user_input,
                        "bot": bot_response,
                        "timestamp": timestamp
                    })
                st.markdown(f"**Bot:** {bot_response}")
                speak_text(bot_response)
            else:
                st.warning("⚠️ Please enter a message.")

    elif input_method == "Voice":
        if st.button("Speak 🎤", key="speak_button_tab_0"):
            user_input = recognize_speech()  
            if user_input:
                with st.spinner("Thinking... 🤔"):
                    bot_response = llama_chatbot(user_input)
                    timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    st.session_state['history'].append({
                        "user": user_input,
                        "bot": bot_response,
                        "timestamp": timestamp
                    })
                st.markdown(f"**Bot:** {bot_response}")
                speak_text(bot_response)

# ---------------------------
# Tab 2: Chat with Documents
# ---------------------------
with tab2:
    st.header("📄 Chat with Documents")

    st.markdown(""" 
    ### 📚 Upload Your Documents
    - Click on the **"Browse files"** button to upload PDF, Word, PPTX, or Excel documents.
    - After uploading, you'll be able to ask questions about the document content using the AI-BOT.
    """)

    uploaded_files = st.file_uploader("Upload Documents", type=["pdf", "docx", "pptx", "xlsx"], accept_multiple_files=True)

    # Check if files are uploaded
    if uploaded_files:
        text = ""
        with st.spinner("Processing documents... 📄"):
            for uploaded_file in uploaded_files:
                if uploaded_file.name.endswith(".pdf"):
                    text += extract_pdf_text([uploaded_file])
                elif uploaded_file.name.endswith(".docx"):
                    text += extract_word_text([uploaded_file])
                elif uploaded_file.name.endswith(".pptx"):
                    text += extract_pptx_text([uploaded_file])
                elif uploaded_file.name.endswith(".xlsx"):
                    text += extract_excel_text([uploaded_file])

        if text:
            st.success("📄 Documents processed successfully!")
            st.text_area("Extracted Text:", text, height=300)

            # Ask questions about the document content
            user_input_doc = st.text_input("Ask questions about the document content:", placeholder="Type your question here...")

            if st.button("Ask 📥", key="document_ask"):
                if user_input_doc.strip():
                    with st.spinner("Thinking... 🤔"):
                        # Using the document's text as context for the chatbot
                        bot_response_doc = llama_chatbot(user_input_doc, context=text)
                        timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                        
                        # Store the conversation history for display
                        if 'history' not in st.session_state:
                            st.session_state['history'] = []
                            
                        st.session_state['history'].append({
                            "user": user_input_doc,
                            "bot": bot_response_doc,
                            "timestamp": timestamp
                        })
                        
                        # Display bot response
                        st.markdown(f"**Bot:** {bot_response_doc}")
                else:
                    st.warning("⚠️ Please enter a question before clicking 'Ask'.")

            # Display chat history related to documents
            if 'history' in st.session_state and st.session_state['history']:
                st.markdown("---")
                st.markdown("### Chat History:")
                for chat in st.session_state['history']:
                    st.markdown(f"**You:** {chat['user']} ({chat['timestamp']})")
                    st.markdown(f"**Bot:** {chat['bot']}")


# ---------------------------
# Tab 3: Chat with Websites
# ---------------------------
with tab3:
    st.header("🌐 Chat with Websites")

    st.markdown(""" 
    ### 🌍 Enter a Website URL
    - Type the URL of the website you want to interact with and hit **Fetch**.
    - After fetching, you can ask questions about the website content using the AI-BOT.
    """)

    # Input field for website URL
    url = st.text_input("Enter the website URL:", placeholder="https://example.com")

    # Initialize session state for website content and chat history
    if 'website_text' not in st.session_state:
        st.session_state['website_text'] = ""
    if 'website_history' not in st.session_state:
        st.session_state['website_history'] = []

    if st.button("Fetch Website Content 🕵️‍♂️", key="fetch_website"):
        if url:
            # Extract website content
            website_text = extract_website_text(url)
            if website_text:
                st.success("🌐 Website content fetched successfully!")
                st.text_area("Website Content:", website_text, height=300)
                
                # Store website text in session state
                st.session_state['website_text'] = website_text

            else:
                st.warning("⚠️ Failed to fetch content from the URL. Please check the URL or try a different one.")
        else:
            st.warning("⚠️ Please enter a URL.")

    # Ask questions about the website content
    if st.session_state['website_text']:
        user_input_web = st.text_input("Ask questions about the website content:", placeholder="Type your question here...")

        if st.button("Ask 📥", key="website_ask"):
            if user_input_web.strip():
                with st.spinner("Thinking... 🤔"):
                    # Using the website content as context for the chatbot
                    bot_response_web = llama_chatbot(user_input_web, context=st.session_state['website_text'])
                    timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    
                    # Store the conversation history for display
                    st.session_state['website_history'].append({
                        "user": user_input_web,
                        "bot": bot_response_web,
                        "timestamp": timestamp
                    })
                    
                    # Display bot response
                    st.markdown(f"**Bot:** {bot_response_web}")
            else:
                st.warning("⚠️ Please enter a question before clicking 'Ask'.")

        # Display chat history related to websites
        if st.session_state['website_history']:
            st.markdown("---")
            st.markdown("### Chat History:")
            for chat in st.session_state['website_history']:
                st.markdown(f"**You:** {chat['user']} ({chat['timestamp']})")
                st.markdown(f"**Bot:** {chat['bot']}")


# ---------------------------
# Tab 4: Chat with Media
# ---------------------------
with tab4:
    st.header("🎥🎤🖼️ Chat with Media")

    st.markdown(""" 
    ### 📹 Video  
    ### 🎧 Audio
    - Use different tabs to upload and interact with video, audio.
    """)

    # Sub-tabs for Chat with Video, Audio, and Image
    media_tabs = st.tabs(["📹 Chat with Video", "🎧 Chat with Audio"])

    # ---------------------------
    # Sub-tab 1: Chat with Video
    # ---------------------------
    with media_tabs[0]:
        st.subheader("📹 Chat with Video")
        uploaded_video = st.file_uploader("Upload Video File", type=["mp4", "avi", "mov"])

        if uploaded_video:
            with st.spinner("Extracting audio... 🎧"):
                audio_file_path = extract_audio_from_video(uploaded_video)
                st.success("🎧 Audio extracted from video successfully!")
                st.audio(audio_file_path, format='audio/aac')

            st.markdown("### Ask questions about the audio content from the video:")
            user_input_video = st.text_input("Your Question (Video):", placeholder="Type your question here...")

            if st.button("Ask Video 📥"):
                if user_input_video.strip():
                    # Placeholder for actual audio transcription logic (optional)
                    st.warning("Audio transcription is not implemented yet.")
                else:
                    st.warning("⚠️ Please enter a question.")

    # ---------------------------
    # Sub-tab 2: Chat with Audio
    # ---------------------------
    with media_tabs[1]:
        st.subheader("🎧 Chat with Audio")
        uploaded_audio = st.file_uploader("Upload Audio File (WAV)", type=["wav"])

        if uploaded_audio:
            with st.spinner("Transcribing audio... 🎙️"):
                audio_text = transcribe_audio(uploaded_audio)
                if audio_text:
                    st.success("🎙️ Audio transcribed successfully!")
                    st.text_area("Transcribed Audio Text:", audio_text, height=200)

                    # Analyze sentiment and emotions for the transcribed audio
                if st.button("Analyze Sentiment and Emotions"):
                    sentiment_result = analyze_sentiment(audio_text)
                    emotion_result = analyze_emotions(audio_text)
                    
                    st.subheader("Sentiment Analysis of Audio:")
                    st.json(sentiment_result)
                    
                    st.subheader("Emotion Analysis of Audio:")
                    st.json(emotion_result)



                    user_input_audio = st.text_input("Ask questions about the audio content:", placeholder="Type your question here...")

                    if st.button("Ask Audio 📥"):
                        if user_input_audio.strip():
                            with st.spinner("Thinking... 🤔"):
                                bot_response_audio = llama_chatbot(user_input_audio, context=audio_text)
                                timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")

                                # Display bot response
                                st.markdown(f"**Bot:** {bot_response_audio}")
                        else:
                            st.warning("⚠️ Please enter a question before clicking 'Ask'.")


# ---------------------------
# Tab 5: Podcast from PDF
# ---------------------------
with tab5:
    st.header("🎙️ Podcast from PDF")

    # Instructions
    st.markdown(""" 
    ### 🎙️ Generate an Audio Podcast
    - Upload a PDF document to convert its content into an audio podcast.
    """)

    uploaded_pdf_podcast = st.file_uploader("Upload PDF Document for Podcast", type=["pdf"])

    if uploaded_pdf_podcast:
        with st.spinner("Extracting text from PDF... 📄"):
            pdf_text = extract_pdf_text([uploaded_pdf_podcast])
            if pdf_text:
                st.success("📄 PDF text extracted successfully!")

                if st.button("Generate Podcast 🎙️"):
                    with st.spinner("Generating podcast... 🎧"):
                        podcast_path = generate_podcast_from_text(pdf_text)
                        if podcast_path:
                            st.audio(podcast_path, format="audio/mp3")
                            st.success("🎙️ Podcast generated successfully!")
                        else:
                            st.error("Failed to generate podcast.")

# ---------------------------
# Tab 6: AI Translator
# ---------------------------
with tab6:
    st.header("🌍 AI Translator")

    # Instructions
    st.markdown(""" 
    ### 🌍 Translate Your Text or Document
    - **Text Input**: Type or paste the text you want to translate.
    - **Document Upload**: Upload a PDF document to translate its content.
    - Select the target language and click **Translate**.
    """)

    # Text input for translation
    translation_input = st.text_area("Enter text to translate:", "")

    # File uploader for PDF documents
    uploaded_pdf_translator = st.file_uploader("Or upload a PDF document:", type=["pdf"])

    # Language selection
    languages = ["en", "hi", "es", "fr", "de", "bn", "ta", "te"]  # Add more languages as needed
    target_language = st.selectbox("Select target language:", languages)

    # Button to perform translation
    if st.button("Translate 🌐"):
        # If the user uploads a PDF, extract text from it
        if uploaded_pdf_translator:
            with st.spinner("Extracting text from PDF... 📄"):
                extracted_text = extract_pdf_text([uploaded_pdf_translator])  # Use your existing function
                translation_input = extracted_text  # Use the extracted text for translation

        # Perform translation if there's text to translate
        if translation_input:
            with st.spinner("Translating... 🌏"):
                translated_text = translate_text(translation_input, target_language)
                st.success("**Translated Text:**")
                st.text_area("", translated_text, height=300)
        else:
            st.warning("⚠️ Please enter text to translate or upload a document.")

# ---------------------------
# Tab 7: Text to SQL
# ---------------------------
with tab7:
    st.header("📊 Text to SQL")

    # Instructions
    st.markdown(""" 
    ### 🔄 Convert Natural Language to SQL
    - Enter your natural language query below to get the corresponding SQL statement.
    - **Example**: "Show me all users from the users table where age is greater than 30."
    """)

    # Text input for SQL query
    user_query = st.text_area("Enter your query:", "")

    if st.button("Convert to SQL 🛠️"):
        if user_query:
            with st.spinner("Converting to SQL... 🛠️"):
                sql_result = text_to_sql(user_query)
                if sql_result:
                    st.success("**SQL Statement:**")
                    st.code(sql_result, language='sql')
                else:
                    st.error("Failed to convert the query to SQL.")
        else:
            st.warning("⚠️ Please enter a query.")

# ---------------------------
# Tab 8: Research Assistant
# ---------------------------
with tab8:
    st.header("🧑‍🏫 Research Assistant")

    st.markdown(""" 
    ### 📚 Upload Research Documents or Input Text
    - You can upload documents or type in text to ask questions about research topics.
    """)

    uploaded_research_files = st.file_uploader("Upload Research Documents (PDF, DOCX)", type=["pdf", "docx"], accept_multiple_files=True)

    if uploaded_research_files:
        research_text = ""
        with st.spinner("Processing research documents... 📄"):
            for uploaded_file in uploaded_research_files:
                if uploaded_file.name.endswith(".pdf"):
                    research_text += extract_pdf_text([uploaded_file])
                elif uploaded_file.name.endswith(".docx"):
                    research_text += extract_word_text([uploaded_file])

        if research_text:
            st.success("📄 Research documents processed successfully!")
            st.text_area("Extracted Research Text:", research_text, height=300)

            user_input_research = st.text_input("Ask questions about your research:", placeholder="Type your question here...")

            if st.button("Ask Research 📥"):
                if user_input_research.strip():
                    with st.spinner("Thinking... 🤔"):
                        bot_response_research = llama_chatbot(user_input_research, context=research_text)
                        timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                        
                        st.session_state['history'].append({
                            "user": user_input_research,
                            "bot": bot_response_research,
                            "timestamp": timestamp
                        })
                        
                        st.markdown(f"**Bot:** {bot_response_research}")
                else:
                    st.warning("⚠️ Please enter a question before clicking 'Ask'.")

            # Display chat history related to research
            if 'history' in st.session_state and st.session_state['history']:
                st.markdown("---")
                st.markdown("### Chat History:")
                for chat in st.session_state['history']:
                    st.markdown(f"**You:** {chat['user']} ({chat['timestamp']})")
                    st.markdown(f"**Bot:** {chat['bot']}")
        else:
            st.warning("⚠️ No text extracted from the uploaded documents. Please check the files.")


# ---------------------------
# Tab: Sentiment and Emotion Analysis
# ---------------------------
with tab9:
    st.header("🧠 Sentiment & Emotion Analysis")

    text_input = st.text_area("Enter Text for Analysis", "")
    
    if text_input:
        sentiment_result = analyze_sentiment(text_input)
        emotion_result = analyze_emotions(text_input)
        
        st.subheader("Sentiment Analysis Result:")
        st.json(sentiment_result)
        
        st.subheader("Emotion Analysis Result:")
        st.json(emotion_result)


